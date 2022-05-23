import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def C_Font (shp,Fname):
    for paragraph in shp.text_frame.paragraphs:
         for run in paragraph.runs:
            run.font.name = Fname

def AutoFont(file,Tpath,Rpath,Fname):

    prs = Presentation(Tpath + '/' + file)

    for slide in prs.slides:
        for shape in slide.shapes:
             if shape.has_text_frame:
                 C_Font(shape,Fname)
             if shape.has_table:
                 for row in shape.table.rows:
                     for cell in row.cells:
                         C_Font(cell,Fname)
             if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                 for shp in shape.shapes:
                     if shp.has_text_frame:
                         C_Font(shp,Fname)
    prs.save(Rpath + '/' + file)


