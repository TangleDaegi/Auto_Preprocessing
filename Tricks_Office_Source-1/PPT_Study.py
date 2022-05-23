from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import openpyxl

def R_TXT (shp,B_W,A_W):
    for paragraph in shp.text_frame.paragraphs:
        for run in paragraph.runs:
            cur_text = run.text
            new_text = cur_text.replace(str(B_W), str(A_W))
            run.text = new_text
def PPT_KC (file,SPath,RPath,DPath):
    wb = openpyxl.load_workbook(DPath)
    ws = wb.active
    before_W = []
    after_W = []
    for cell in ws.rows:
        before_W.append(cell[0].value)
        after_W.append(cell[1].value)

    prs = Presentation(SPath + "/" + file)

    for slide in prs.slides:
        for shape in slide.shapes:
            for i in range(1, len(after_W)):
                if shape.has_text_frame:
                    R_TXT(shape,before_W[i],after_W[i])
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            R_TXT(cell,before_W[i],after_W[i])
                if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                    for shp in shape.shapes:
                        if shp.has_text_frame:
                            R_TXT(shp,before_W[i],after_W[i])
    prs.save(RPath + "/" + file)