from googletrans import Translator
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def Trans(shape, language, trans):
    text_frame = shape.text_frame
    for paragraph in text_frame.paragraphs:
        try:
            font = paragraph.runs[0].font
            Fname = font.name
            Fsize = font.size
            Fbold = font.bold
        except:
            pass
        cur_text = paragraph.text
        new_text = trans.translate(cur_text, dest=language).text
        paragraph.text = new_text
        for run in paragraph.runs:
            font = run.font
            font.name = Fname
            font.size = Fsize
            font.bold = Fbold

def GTrans(file,TPath,RPath,language):
    trans=Translator()
    prs=Presentation(TPath + "/" + file)
    Slides = prs.slides
    for slide in Slides:
        Shapes = slide.shapes
        for shape in Shapes:
            if shape.has_text_frame:
                Trans(shape, language,trans)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        Trans(cell, language,trans)
            if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                for shp in shape.shapes:
                    if shp.has_text_frame:
                        Trans(shp, language,trans)
    prs.save(RPath + "/" + file)